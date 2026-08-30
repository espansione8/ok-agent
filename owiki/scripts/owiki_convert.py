#!/usr/bin/env python3
"""owiki_convert.py — RAW→WIKI converter for Obsidian vaults.

Converts <vault>/<project>/RAW/ files into Obsidian markdown in
<vault>/<project>/WIKI/ as project knowledge for coding agents.
Vault resolution: --vault > OWIKI_VAULT_PATH > vault-config.json >
auto-detection (agent mode) / interactive wizard (terminal) > hard-gate
error. Never guessed from cwd. --vault and auto-detection SAVE the
config, so one successful run configures every future run.

Usage:
  owiki_convert.py <project_dir> [--force]   convert one folder (WIKI only)
  owiki_convert.py <name> [--force]          named vault subfolder (only form
                                             writing KNOWLEDGE.md/AGENTS.md)
  owiki_convert.py --all | --list | --init | --update

No args: cwd == vault root → all subfolders; cwd in a subfolder → just it;
cwd outside → notify only. No-arg modes create RAW/ (moving loose supported
files in) and never write agent files.

Cache: <project>/.owiki-cache.json (SHA-256 + size). Unchanged files skip,
deleted sources prune, WIKI mirrors RAW exactly.

Figures (PDF raster/vector, standalone images) → WIKI/assets/ with a
mechanical caption + Obsidian embed. Semantic descriptions and English
translation are enrichment handled by the active owiki agent model, not
by a separate service or model setting.

Requires: beautifulsoup4 markdownify pymupdf python-docx openpyxl
python-pptx pillow pyyaml; LibreOffice only for ODF/legacy.
"""
import sys, os, re, json, hashlib, argparse, subprocess, tempfile, time
from datetime import datetime

try:
    import yaml
except ImportError:
    print("ERROR: PyYAML is required (pip install pyyaml).")
    sys.exit(1)

if hasattr(sys.stdout, 'reconfigure'):
    try:
        sys.stdout.reconfigure(encoding='utf-8', errors='replace')
    except Exception:
        pass

SOFFICE_CANDIDATES = [
    r'C:\Program Files\LibreOffice\program\soffice.exe',
    r'C:\Program Files (x86)\LibreOffice\program\soffice.exe',
    '/usr/bin/soffice', '/usr/local/bin/soffice', '/opt/libreoffice/program/soffice',
]


def find_soffice():
    for cand in SOFFICE_CANDIDATES:
        if os.path.isfile(cand):
            return cand
    from shutil import which
    return which('soffice') or which('libreoffice')


SOFFICE = find_soffice()

_CTX = {'assets_dir': None, 'slug': None, 'assets': []}


def _asset_rel(name):
    return f'assets/{name}'


_DEPS = {}


def _get_dep(name):
    if name in _DEPS:
        return _DEPS[name]
    try:
        if name == 'bs4':
            from bs4 import BeautifulSoup
            _DEPS[name] = BeautifulSoup
        elif name == 'markdownify':
            from markdownify import markdownify
            _DEPS[name] = markdownify
        elif name == 'fitz':
            import pymupdf as fitz
            _DEPS[name] = fitz
        elif name == 'docx':
            from docx import Document
            from docx.oxml.text.paragraph import CT_P
            from docx.oxml.table import CT_Tbl
            from docx.text.paragraph import Paragraph
            from docx.table import Table
            _DEPS[name] = (Document, CT_P, CT_Tbl, Paragraph, Table)
        elif name == 'openpyxl':
            from openpyxl import load_workbook
            _DEPS[name] = load_workbook
        elif name == 'pptx':
            from pptx import Presentation
            from pptx.enum.shapes import PP_PLACEHOLDER
            _DEPS[name] = (Presentation, PP_PLACEHOLDER)
        elif name == 'PIL':
            from PIL import Image
            _DEPS[name] = Image
        return _DEPS[name]
    except ImportError:
        _DEPS[name] = None
        return None


def file_hash(filepath):
    h = hashlib.sha256()
    size = 0
    with open(filepath, 'rb') as f:
        while True:
            chunk = f.read(8192)
            if not chunk:
                break
            h.update(chunk)
            size += len(chunk)
    return h.hexdigest(), size


def load_cache(cache_path):
    if os.path.isfile(cache_path):
        try:
            with open(cache_path, 'r', encoding='utf-8') as f:
                return json.load(f)
        except (json.JSONDecodeError, IOError):
            pass
    return {}


def save_cache(cache_path, cache):
    with open(cache_path, 'w', encoding='utf-8') as f:
        json.dump(cache, f, indent=2, ensure_ascii=False)


def clean_markdown(text):
    text = re.sub(r'\n{3,}', '\n\n', text)
    return '\n'.join(line.rstrip() for line in text.split('\n')).strip()


def summarize(text, max_len=150):
    # 1. Try the first markdown heading content (usually a section title)
    for line in text.split('\n'):
        stripped = line.strip()
        if stripped.startswith('#') and not stripped.startswith('# '):
            content = stripped.lstrip('#').strip()
            if (len(content) > 10
                    and not content.lower().startswith('figure ')
                    and not content.lower().startswith('include ')
                    and not content.lower().startswith('define ')):
                return content[:max_len]
    # 2. First non-trivial, non-embed, non-header line
    for line in text.split('\n'):
        stripped = line.strip()
        if not stripped:
            continue
        if stripped.startswith(('#', '|', '---', '>', '![[', '```', '<!--')):
            continue
        if len(stripped) > 30:
            return stripped if len(stripped) <= max_len else stripped[:max_len - 3] + '...'
    # 3. Fallback — empty is better than fake text (enrichment fills it)
    return ''


def _sniff_html_encoding(raw):
    if raw.startswith(b'\xef\xbb\xbf'):
        return 'utf-8-sig'
    if raw.startswith(b'\xff\xfe'):
        return 'utf-16-le'
    if raw.startswith(b'\xfe\xff'):
        return 'utf-16-be'
    head = raw[:4096].decode('ascii', errors='ignore').lower()
    m = re.search(r'charset=["\']?\s*([\w\-]+)', head)
    return m.group(1) if m else 'utf-8'


def convert_html(filepath):
    BeautifulSoup = _get_dep('bs4')
    markdownify = _get_dep('markdownify')
    if BeautifulSoup is None or markdownify is None:
        return None, ''
    with open(filepath, 'rb') as f:
        raw = f.read()
    try:
        html = raw.decode(_sniff_html_encoding(raw), errors='replace')
    except LookupError:
        html = raw.decode('utf-8', errors='replace')
    soup = BeautifulSoup(html, 'html.parser')
    for tag in soup.find_all(['style', 'script', 'meta', 'link']):
        tag.decompose()
    # Demote h1 → h2 in source HTML: the note template already provides # {title}
    for h1 in soup.find_all('h1'):
        h1.name = 'h2'
    title_tag = soup.find('title')
    title = title_tag.get_text().strip() if title_tag else ''
    body = soup.find('body') or soup
    md = markdownify(str(body), heading_style='ATX',
                     strip=['script', 'meta', 'link'])
    return clean_markdown(md), title


def convert_txt(filepath):
    with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
        content = f.read()
    content = content.replace('\r\n', '\n').replace('\r', '\n')
    if filepath.lower().endswith('.md'):  # strip source's own frontmatter
        fm_match = re.match(r'^---\n.*?\n---\n?', content, re.DOTALL)
        if fm_match:
            content = content[fm_match.end():]
    return clean_markdown(content), ''


MIN_IMAGE_PX = 200
VECTOR_DRAWINGS_THRESHOLD = 10
VECTOR_TEXT_LEN_MAX = 200  # render only text-sparse pages (real schematics)
RENDER_DPI = 150


def _save_pixmap(pix, dest_path):
    fitz = _get_dep('fitz')
    if pix.colorspace and pix.colorspace.name not in ('DeviceRGB', 'DeviceGray'):
        pix = fitz.Pixmap(fitz.csRGB, pix)
    pix.save(dest_path)


def _caption_for(fallback_info):
    """Return a deterministic caption; semantic captioning belongs to the active agent."""
    return fallback_info


def _strip_running_headers(pages_text):
    """Remove lines that repeat on most pages (running headers/footers).

    Detects lines that appear on >50% of pages with near-identical content
    (digits normalised to #) and strips them from every page.
    """
    if len(pages_text) < 3:
        return pages_text
    from collections import Counter
    normalized = []
    for page_text in pages_text:
        page_lines = [re.sub(r'\d+', '#', line.strip()) for line in page_text.split('\n')]
        normalized.append(page_lines)
    line_pages = Counter()
    for page_lines in normalized:
        for line in set(page_lines):
            if len(line) > 3:
                line_pages[line] += 1
    threshold = max(3, len(pages_text) * 0.5)
    noise_lines = {l for l, c in line_pages.items() if c >= threshold}
    if not noise_lines:
        return pages_text
    result = []
    for page_text in pages_text:
        filtered = [line for line in page_text.split('\n')
                    if re.sub(r'\d+', '#', line.strip()) not in noise_lines]
        result.append('\n'.join(filtered))
    return result


def _detect_lang(block):
    """Heuristic language detection for fenced code blocks."""
    joined = '\n'.join(block)
    if '<?xml' in joined or re.match(r'^\s*<\w+', joined):
        return 'xml'
    if re.search(r'#include\b|int\s+main\s*\(', joined):
        return 'cpp'
    if re.search(r'^\s*def\s+\w+|^import\s+\w+|from\s+\w+\s+import', joined, re.MULTILINE):
        return 'python'
    if joined.strip().startswith('{') and '"key"' not in joined and ':' in joined:
        # Could be JSON or a dict literal; check for quoted keys
        if re.search(r'["\']\w+["\']\s*:', joined):
            return 'json'
    if joined.strip().startswith('{') and re.search(r'["\']\w+["\']\s*:', joined):
        return 'json'
    return ''


def _fence_code_blocks(text):
    """Detect line-numbered code sequences and fence them.

    PDF text extraction often emits code with all the line numbers grouped
    together (1\\n2\\n3\\n...\\nN\\n) followed by the code content lines, OR
    interleaved (1\\n<code>\\n2\\n<code>\\n...). This detects both patterns,
    strips the numbers, and wraps the code in ```lang fences.
    """
    lines = text.split('\n')
    result = []
    i = 0
    while i < len(lines):
        stripped = lines[i].strip()
        # Bare number on its own line → potential code-block line number
        if re.match(r'^\d+$', stripped):
            # Collect the full sequence of consecutive bare numbers
            numbers = []
            j = i
            expected = int(stripped)
            while j < len(lines) and re.match(rf'^{expected}$', lines[j].strip()):
                numbers.append(expected)
                j += 1
                expected += 1
            if len(numbers) >= 5:
                # Pattern A: all numbers grouped, then code block follows
                # Check if the line after the number sequence is non-empty
                # and looks like code (not another heading or paragraph)
                if j < len(lines):
                    code_lines = []
                    k = j
                    # Collect code until we hit a blank line followed by non-code,
                    # a page comment, a heading, or end of file
                    while k < len(lines):
                        line = lines[k]
                        next_stripped = line.strip()
                        if not next_stripped:
                            # Allow blank lines in code — only stop if we see
                            # 3+ consecutive blanks (definitely end of code block)
                            blk_count = 0
                            for m in range(k, min(k + 3, len(lines))):
                                if not lines[m].strip():
                                    blk_count += 1
                                else:
                                    break
                            if blk_count >= 3:
                                break
                            code_lines.append(line)
                            k += 1
                            continue
                        if next_stripped.startswith('<!-- page'):
                            break
                        # Stop at markdown headings (# followed by space) but not
                        # C preprocessor directives (#include, #define, etc.)
                        if next_stripped.startswith('#') and (
                            len(next_stripped) == 1 or next_stripped[1] == ' '
                        ):
                            break
                        if next_stripped.startswith('![[assets/'):
                            break
                        if next_stripped.startswith('### Figure'):
                            break
                        # Stop if we see a new number sequence starting (1, 2...)
                        # that belongs to a different code block
                        if re.match(r'^1$', next_stripped):
                            if k + 1 < len(lines) and re.match(r'^2$', lines[k + 1].strip()):
                                break
                        # Stop at prose lines that introduce a new code block
                        # (contain words and end with ':')
                        if (len(next_stripped) > 15 and next_stripped.endswith(':')
                                and ' ' in next_stripped
                                and not next_stripped.startswith('//')
                                and not next_stripped.startswith('#')):
                            break
                        code_lines.append(line)
                        k += 1
                    if len(code_lines) >= 3:
                        lang = _detect_lang(code_lines)
                        result.append(f'```{lang}')
                        result.extend(code_lines)
                        result.append('```')
                        i = k
                        continue
                # Pattern B: interleaved (number → code → number → code)
                # Not seen in practice but handle gracefully
                block = []
                j2 = i
                exp2 = int(stripped)
                while j2 < len(lines) and re.match(rf'^{exp2}$', lines[j2].strip()):
                    if j2 + 1 < len(lines):
                        next_line = lines[j2 + 1].strip()
                        if re.match(r'^\d+$', next_line):
                            break
                        block.append(lines[j2 + 1])
                        j2 += 2
                        exp2 += 1
                    else:
                        break
                if len(block) >= 3:
                    lang = _detect_lang(block)
                    result.append(f'```{lang}')
                    result.extend(block)
                    result.append('```')
                    i = j2
                    continue
        result.append(lines[i])
        i += 1
    return '\n'.join(result)


def convert_pdf(filepath):
    fitz = _get_dep('fitz')
    if fitz is None:
        return None, ''
    doc = fitz.open(filepath)
    md_lines = []
    assets_dir = _CTX['assets_dir']
    slug = _CTX['slug']
    _CTX['assets'] = []
    seen_xrefs = set()
    figure_no = 0
    pages_text = []
    pages_images = []
    for page_num, page in enumerate(doc):
        text = page.get_text("text")
        image_sections = []
        for img in page.get_images(full=True):
            xref = img[0]
            if xref in seen_xrefs:
                continue
            try:
                pix = fitz.Pixmap(doc, xref)
            except Exception:
                continue
            w, h = pix.width, pix.height
            if w < MIN_IMAGE_PX or h < MIN_IMAGE_PX:
                continue
            seen_xrefs.add(xref)
            figure_no += 1
            name = f'{slug}-p{page_num + 1}-fig{figure_no}.png'
            dest = os.path.join(assets_dir, name)
            try:
                _save_pixmap(pix, dest)
                _CTX['assets'].append(_asset_rel(name))
            except Exception:
                continue
            caption = _caption_for(
                f'Embedded figure {figure_no} from page {page_num + 1} '
                f'({w}×{h}px). The active owiki model can add a semantic description.')
            image_sections.append(
                f'### Figure {figure_no} — page {page_num + 1}\n'
                f'> {caption}\n'
                f'![[{_asset_rel(name)}]]')
            print(f"      🖼  {name} ({w}×{h}px)")
        drawings = page.get_drawings()
        if len(drawings) > VECTOR_DRAWINGS_THRESHOLD and len(text.strip()) < VECTOR_TEXT_LEN_MAX:
            figure_no += 1
            name = f'{slug}-p{page_num + 1}-fig{figure_no}.png'
            dest = os.path.join(assets_dir, name)
            try:
                pix = page.get_pixmap(dpi=RENDER_DPI)
                _save_pixmap(pix, dest)
                _CTX['assets'].append(_asset_rel(name))
            except Exception:
                continue
            caption = _caption_for(
                f'Vector drawing {figure_no} from page {page_num + 1} '
                f'(rendered at {RENDER_DPI} DPI). The active owiki model can add a semantic description.')
            image_sections.append(
                f'### Figure {figure_no} — page {page_num + 1} (schematic)\n'
                f'> {caption}\n'
                f'![[{_asset_rel(name)}]]')
            print(f"      🖼  {name} (vector render)")
        pages_text.append(text)
        pages_images.append(image_sections)
    doc.close()

    # Strip repeated running headers/footers across pages
    pages_text = _strip_running_headers(pages_text)

    # Assemble — HTML comments for page provenance, no ## Page headers, no --- separators
    for idx, (text, image_sections) in enumerate(zip(pages_text, pages_images)):
        page_md = []
        if text.strip():
            page_md.append(f'<!-- page {idx + 1} -->\n{text}')
        if image_sections:
            page_md.append('\n'.join(image_sections))
        if page_md:
            md_lines.append('\n'.join(page_md))

    body = '\n\n'.join(md_lines)
    # Fence code blocks (detect line-number prefixes, strip them, wrap in ```lang)
    body = _fence_code_blocks(body)
    return clean_markdown(body), os.path.splitext(os.path.basename(filepath))[0]


def convert_image(filepath):
    Image = _get_dep('PIL')
    if Image is None:
        return None, ''
    slug = _CTX['slug']
    assets_dir = _CTX['assets_dir']
    _CTX['assets'] = []
    ext = os.path.splitext(filepath)[1].lower()
    name = f'{slug}{ext}'
    dest = os.path.join(assets_dir, name)
    import shutil
    shutil.copy2(filepath, dest)
    _CTX['assets'].append(_asset_rel(name))
    try:
        with Image.open(filepath) as im:
            w, h = im.size
    except Exception:
        w = h = 0
    caption = _caption_for(
        f'Image ({w}×{h}px). The active owiki model can add a semantic description.')
    md = (f'## Image\n> {caption}\n\n'
          f'![[{_asset_rel(name)}]]\n\n'
          f'Source: RAW/{os.path.basename(filepath)}')
    print(f"      🖼  {name} ({w}×{h}px)")
    return clean_markdown(md), slug.replace('-', ' ')


def convert_docx(filepath):
    deps = _get_dep('docx')
    if deps is None:
        return None, ''
    Document, CT_P, CT_Tbl, Paragraph, Table = deps
    doc = Document(filepath)
    md_lines = []
    for child in doc.element.body.iterchildren():
        if isinstance(child, CT_P):
            para = Paragraph(child, doc)
            style = (para.style.name if para.style else '') or ''
            text = para.text.strip()
            if not text:
                continue
            if 'Title' in style:
                md_lines.append(f'# {text}')
            elif 'Heading 1' in style:
                md_lines.append(f'## {text}')
            elif 'Heading 2' in style:
                md_lines.append(f'### {text}')
            elif 'Heading 3' in style:
                md_lines.append(f'#### {text}')
            elif 'List Bullet' in style or 'List' in style:
                md_lines.append(f'- {text}')
            elif 'List Number' in style:
                md_lines.append(f'1. {text}')
            else:
                md_lines.append(text)
        elif isinstance(child, CT_Tbl):
            table = Table(child, doc)
            rows = []
            for row in table.rows:
                # dedupe horizontally-merged cells (shared <w:tc>)
                cells = []
                prev_tc = None
                for c in row.cells:
                    if prev_tc is not None and c._tc is prev_tc:
                        cells.append('')
                    else:
                        cells.append(c.text.strip().replace('|', '\\|').replace('\n', ' '))
                    prev_tc = c._tc
                rows.append('| ' + ' | '.join(cells) + ' |')
            if rows:
                ncols = len(rows[0].split('|')) - 2
                sep = '| ' + ' | '.join(['---'] * ncols) + ' |'
                md_lines.append('\n' + '\n'.join([rows[0], sep] + rows[1:]) + '\n')
    return clean_markdown('\n'.join(md_lines)), ''


def convert_xlsx(filepath):
    load_workbook = _get_dep('openpyxl')
    if load_workbook is None:
        return None, ''
    wb = load_workbook(filepath, data_only=True)
    md_lines = []
    for sheet_name in wb.sheetnames:
        ws = wb[sheet_name]
        rows_data = []
        for row in ws.iter_rows(values_only=True):
            if all(v is None or str(v).strip() == '' for v in row):
                continue
            rows_data.append(['' if v is None else str(v).replace('|', '\\|') for v in row])
        if not rows_data:
            continue
        max_cols = max(len(r) for r in rows_data)
        rows_data = [r + [''] * (max_cols - len(r)) for r in rows_data]
        while max_cols > 1 and all(not r[max_cols - 1].strip() for r in rows_data):
            max_cols -= 1
        rows_data = [r[:max_cols] for r in rows_data]
        md_lines.append(f'## Sheet: {sheet_name}\n')
        header = rows_data[0]
        md_lines.append('| ' + ' | '.join(header) + ' |')
        md_lines.append('| ' + ' | '.join(['---'] * len(header)) + ' |')
        for r in rows_data[1:]:
            md_lines.append('| ' + ' | '.join(r) + ' |')
        md_lines.append('')
    return clean_markdown('\n'.join(md_lines)), ''


def convert_pptx(filepath):
    deps = _get_dep('pptx')
    if deps is None:
        return None, ''
    Presentation, PP_PLACEHOLDER = deps
    prs = Presentation(filepath)
    md_lines = []
    for i, slide in enumerate(prs.slides, 1):
        md_lines.append(f'## Slide {i}\n')
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            is_title = False
            try:
                if shape.is_placeholder:
                    is_title = shape.placeholder_format.type in (
                        PP_PLACEHOLDER.TITLE, PP_PLACEHOLDER.CENTER_TITLE)
            except Exception:
                pass
            first = True
            for para in shape.text_frame.paragraphs:
                text = ''.join(run.text for run in para.runs).strip()
                if not text:
                    continue
                if is_title and first:
                    md_lines.append(f'### {text}')
                elif is_title:
                    md_lines.append(text)
                elif para.level == 0:
                    md_lines.append(f'- {text}')
                else:
                    indent = '  ' * (para.level - 1)
                    md_lines.append(f'  {indent}- {text}')
                first = False
    return clean_markdown('\n'.join(md_lines)), ''


def convert_csv(filepath):
    with open(filepath, 'r', encoding='utf-8', errors='replace') as f:
        import csv as csv_mod
        rows = [list(r) for r in csv_mod.reader(f) if any(c.strip() for c in r)]
    if not rows:
        return '', ''
    max_cols = max(len(r) for r in rows)
    rows = [r + [''] * (max_cols - len(r)) for r in rows]
    md_lines = ['| ' + ' | '.join(rows[0]) + ' |',
                '| ' + ' | '.join(['---'] * max_cols) + ' |']
    for r in rows[1:]:
        md_lines.append('| ' + ' | '.join(r) + ' |')
    return clean_markdown('\n'.join(md_lines)), ''


LO_MAP = {
    '.odt': ('docx', convert_docx), '.doc': ('docx', convert_docx),
    '.rtf': ('docx', convert_docx), '.ods': ('xlsx', convert_xlsx),
    '.xls': ('xlsx', convert_xlsx), '.odp': ('pptx', convert_pptx),
    '.ppt': ('pptx', convert_pptx),
}


def soffice_convert(src_path, target_ext, dest_dir, timeout=180):
    global SOFFICE
    if SOFFICE is None:
        SOFFICE = find_soffice()
    if SOFFICE is None:
        return None, "LibreOffice not found (needed for this format)"
    os.makedirs(dest_dir, exist_ok=True)
    profile = '-env:UserInstallation=file:///' + \
        tempfile.gettempdir().replace('\\', '/') + '/lo_profile_owiki'
    result = None
    last_err = None
    for attempt in range(3):
        try:
            result = subprocess.run(
                [SOFFICE, '--headless', profile,
                 '--convert-to', target_ext, '--outdir', dest_dir, src_path],
                capture_output=True, text=True, timeout=timeout)
        except subprocess.TimeoutExpired:
            last_err = f'soffice timed out after {timeout}s'
            time.sleep(3)
            continue
        except (PermissionError, OSError) as e:
            last_err = f'soffice launch failed: {e}'
            time.sleep(3)
            continue
        matches = [f for f in os.listdir(dest_dir) if f.lower().endswith('.' + target_ext)]
        if matches:
            return os.path.join(dest_dir, matches[0]), None
        last_err = f"soffice produced no output (rc={result.returncode}, err={result.stderr[:200]})"
        if attempt < 2:
            time.sleep(2)
    return None, last_err or "soffice failed for an unknown reason after 3 attempts"


def convert_lo_format(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    target_ext, native_conv = LO_MAP[ext]
    with tempfile.TemporaryDirectory(prefix='owiki_lo_') as tmp:
        inter, err = soffice_convert(filepath, target_ext, tmp)
        if inter is None:
            return None, err
        return native_conv(inter)


def slugify(filename):
    name = os.path.splitext(filename)[0]
    name = re.sub(r'[_\s]+', '-', name)
    name = re.sub(r'[^\w\-.]', '', name)
    return '-'.join(p.capitalize() for p in name.split('-') if p) or 'Untitled'


def infer_category(filename):
    ext = os.path.splitext(filename)[1].lower()
    if ext == '.pdf':
        return 'reference'
    if ext in ('.txt', '.md'):
        return 'notes'
    if ext in ('.html', '.htm', '.docx', '.doc', '.odt', '.rtf'):
        return 'document'
    if ext in ('.xlsx', '.xls', '.ods', '.csv'):
        return 'data'
    if ext in ('.pptx', '.ppt', '.odp'):
        return 'presentation'
    if ext in ('.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp', '.tiff', '.tif'):
        return 'image'
    return 'general'


# Generic filler-word filter — no project/domain vocabulary.
_TAG_STOPWORDS = {
    'the', 'and', 'for', 'with', 'from', 'this', 'that', 'of', 'to', 'in',
    'on', 'at', 'by', 'is', 'are', 'was', 'were', 'ver', 'version', 'final',
    'draft', 'copy', 'new', 'old', 'untitled',
}
MAX_INFERRED_TAGS = 6


def infer_tags(filename, title):
    stem = os.path.splitext(filename)[0]
    text = f'{stem} {title}'.lower()
    words = re.findall(r'[^\W\d_]{3,}', text, re.UNICODE)
    tags = []
    for w in words:
        if w in _TAG_STOPWORDS or w in tags:
            continue
        tags.append(w)
        if len(tags) >= MAX_INFERRED_TAGS:
            break
    return tags or ['general']


RELATED_MARKER = '<!-- owiki:related -->'


def build_related_links(current_slug, all_notes):
    if current_slug not in all_notes:
        return ''
    current_tags = set(all_notes[current_slug]['tags'])
    links = []
    for slug, data in all_notes.items():
        if slug != current_slug and current_tags & set(data['tags']):
            links.append(f'- [[{slug}|{data["title"]}]]')
    links.append('- [[_Index|← Back to Index]]')
    body = '\n'.join(sorted(set(links))) if len(links) > 1 else '- [[_Index|← Back to Index]]'
    return f'\n{RELATED_MARKER}\n## Related\n{body}'


def build_index(project_name, all_notes):
    now = datetime.now().strftime('%Y-%m-%d')
    categories = {}
    for slug, data in all_notes.items():
        categories.setdefault(data['category'], []).append(data)
    lines = ['---', f'title: "{project_name} — Wiki Index"', f'converted: {now}',
             '---', '', f'# {project_name} — Wiki Index', '',
             f'> **{len(all_notes)} notes** across **{len(categories)} categories**',
             f'> — last updated {now}', '', '---', '']
    cat_icons = {'document': '📄', 'reference': '📚', 'notes': '📝',
                 'data': '📊', 'presentation': '🖼️', 'image': '🖼️',
                 'general': '📁'}
    for cat in sorted(categories):
        lines.append(f'## {cat_icons.get(cat, "📁")} {cat.replace("-", " ").title()}')
        lines += ['', '| Note | Summary |', '| --- | --- |']
        for item in sorted(categories[cat], key=lambda x: x['title']):
            link = f'[[{item["slug"]}|{item["title"]}]]'
            summary = item['summary'].replace('|', '\\|')
            lines.append(f'| {link} | {summary} |')
        lines.append('')
    lines += ['---', '', '## How to Use', '',
              '- Each note has YAML frontmatter with `title`, `category`, `tags`,',
              '  and `source` fields for programmatic access.',
              '- Cross-links use Obsidian `[[wikilink]]` syntax.',
              '- Original source is preserved in `source` frontmatter field.', '']
    return '\n'.join(lines)


def build_frontmatter(title, source_filename, category, tags, project_name, summary, images=None):
    meta = {'title': title, 'aliases': [title, source_filename],
            'category': category, 'source': f'RAW/{source_filename}',
            'project': project_name, 'converted': datetime.now().strftime('%Y-%m-%d'),
            'tags': tags, 'summary': summary, 'images': images or []}
    dumped = yaml.safe_dump(meta, sort_keys=False, allow_unicode=True,
                            default_flow_style=False).strip()
    return f"---\n{dumped}\n---"


def parse_note_frontmatter(content):
    fm_match = re.match(r'^---\n(.*?)\n---', content, re.DOTALL)
    if not fm_match:
        return None
    try:
        meta = yaml.safe_load(fm_match.group(1))
    except yaml.YAMLError:
        return None
    return meta if isinstance(meta, dict) else None


def _metadata_list(value):
    if isinstance(value, list):
        return [str(item) for item in value if item]
    return [str(value)] if value else []


def _asset_basename(asset_path):
    """Return the generated asset filename, never allowing path traversal."""
    return os.path.basename(os.path.normpath(str(asset_path)))


def cleanup_generated_outputs(wiki_dir, assets_dir, current_wiki_files, current_assets):
    """Remove generated notes/assets that no longer have a RAW source."""
    removed_notes = 0
    for name in os.listdir(wiki_dir):
        path = os.path.join(wiki_dir, name)
        if (os.path.isfile(path) and name.lower().endswith('.md')
                and name not in current_wiki_files):
            os.remove(path)
            print(f"    ✓ Removed orphan: {name}")
            removed_notes += 1

    current_assets = {_asset_basename(name) for name in current_assets}
    removed_assets = 0
    if os.path.isdir(assets_dir):
        for root, dirs, files in os.walk(assets_dir, topdown=False):
            for name in files:
                path = os.path.join(root, name)
                if os.path.relpath(path, assets_dir) not in current_assets:
                    os.remove(path)
                    removed_assets += 1
            for name in dirs:
                path = os.path.join(root, name)
                try:
                    os.rmdir(path)
                except OSError:
                    pass
    if removed_assets:
        print(f"    ✓ Removed {removed_assets} orphaned asset(s)")
    return removed_notes, removed_assets


SUPPORTED_EXT = {'.html', '.htm', '.txt', '.md', '.pdf',
                 '.docx', '.xlsx', '.pptx', '.csv',
                 '.odt', '.ods', '.odp', '.doc', '.xls', '.ppt', '.rtf',
                 '.jpg', '.jpeg', '.png', '.gif', '.webp', '.bmp', '.tiff', '.tif'}

NATIVE = {'.html': convert_html, '.htm': convert_html, '.txt': convert_txt,
          '.md': convert_txt, '.pdf': convert_pdf, '.docx': convert_docx,
          '.xlsx': convert_xlsx, '.pptx': convert_pptx, '.csv': convert_csv,
          '.jpg': convert_image, '.jpeg': convert_image, '.png': convert_image,
          '.gif': convert_image, '.webp': convert_image, '.bmp': convert_image,
          '.tiff': convert_image, '.tif': convert_image}


def convert_file(filepath):
    ext = os.path.splitext(filepath)[1].lower()
    if ext in NATIVE:
        return NATIVE[ext](filepath)
    if ext in LO_MAP:
        return convert_lo_format(filepath)
    return None, ''


def run_project(project_dir, force):
    project_dir = os.path.abspath(project_dir)
    raw_dir = os.path.join(project_dir, 'RAW')
    wiki_dir = os.path.join(project_dir, 'WIKI')
    assets_dir = os.path.join(wiki_dir, 'assets')
    cache_path = os.path.join(project_dir, '.owiki-cache.json')
    if not os.path.isdir(raw_dir):
        print(f"ERROR: RAW directory not found: {raw_dir}")
        return 1
    os.makedirs(wiki_dir, exist_ok=True)
    os.makedirs(assets_dir, exist_ok=True)
    project_name = os.path.basename(project_dir)
    cache = load_cache(cache_path)
    cache_changed = False

    raw_files = sorted([f for f in os.listdir(raw_dir)
                        if os.path.splitext(f)[1].lower() in SUPPORTED_EXT])
    raw_files_set = set(raw_files)
    if not raw_files:
        print(f"No supported files in {raw_dir}")
        print(f"Supported: {', '.join(sorted(SUPPORTED_EXT))}")
        # fall through: empty RAW must still orphan-clean WIKI

    notes_data = {}
    to_convert = []
    slug_counts = {}
    file_slugs = {}
    for fname in raw_files:
        base_slug = slugify(fname)
        count = slug_counts.get(base_slug, 0)
        slug_counts[base_slug] = count + 1
        file_slugs[fname] = base_slug if count == 0 else f'{base_slug}-{count}'

    for fname in raw_files:
        fpath = os.path.join(raw_dir, fname)
        fhash, fsize = file_hash(fpath)
        slug = file_slugs[fname]
        wiki_filename = slug + '.md'
        wiki_path = os.path.join(wiki_dir, wiki_filename)
        cached = cache.get(fname)
        if cached and not force:
            if (cached.get('hash') == fhash and cached.get('size') == fsize
                    and cached.get('wiki_file') == wiki_filename
                    and os.path.isfile(wiki_path)):
                with open(wiki_path, 'r', encoding='utf-8') as f:
                    meta = parse_note_frontmatter(f.read())
                if meta is not None:
                    tags = meta.get('tags') or ['general']
                    if not isinstance(tags, list):
                        tags = [str(tags)]
                    notes_data[slug] = {
                        'slug': slug, 'title': str(meta.get('title', slug)),
                        'category': str(meta.get('category', 'general')),
                        'tags': [str(t) for t in tags],
                        'summary': str(meta.get('summary', '')),
                        'assets': _metadata_list(meta.get('images'))}
                    print(f"  ✓ Skip (unchanged): {fname}")
                    continue
                # unparseable frontmatter → reconvert, never silently drop
        to_convert.append((fname, slug, fpath, fhash, fsize))

    converted_count = 0
    failed_count = 0
    for fname, slug, fpath, fhash, fsize in to_convert:
        print(f"  Converting: {fname} → {slug}.md")
        _CTX['assets_dir'] = assets_dir
        _CTX['slug'] = slug
        _CTX['assets'] = []
        try:
            body_md, extracted_title = convert_file(fpath)
        except Exception as exc:
            body_md, extracted_title = None, str(exc)
        file_assets = list(_CTX['assets'])
        if body_md is None:
            failed_count += 1
            detail = f": {extracted_title}" if extracted_title else ''
            print(f"    ⚠ Skipped (conversion failed){detail}: {fname}")
            for asset_rel in file_assets:
                asset_path = os.path.join(assets_dir, _asset_basename(asset_rel))
                if os.path.isfile(asset_path):
                    os.remove(asset_path)
            # keep last good note alive so orphan cleanup can't delete it
            old_entry = cache.get(fname)
            old_wiki_file = old_entry.get('wiki_file') if old_entry else None
            if old_wiki_file:
                old_path = os.path.join(wiki_dir, old_wiki_file)
                if os.path.isfile(old_path):
                    with open(old_path, 'r', encoding='utf-8') as f:
                        old_meta = parse_note_frontmatter(f.read())
                    if old_meta is not None:
                        old_slug = os.path.splitext(old_wiki_file)[0]
                        old_tags = old_meta.get('tags') or ['general']
                        notes_data[old_slug] = {
                            'slug': old_slug,
                            'title': str(old_meta.get('title', old_slug)),
                            'category': str(old_meta.get('category', 'general')),
                            'tags': [str(t) for t in (old_tags if isinstance(old_tags, list) else [old_tags])],
                            'summary': str(old_meta.get('summary', '')),
                            'assets': _metadata_list(old_meta.get('images'))}
                        print(f"    ↳ keeping last good version: {old_wiki_file}")
            continue
        title = extracted_title or slug.replace('-', ' ')
        category = infer_category(fname)
        tags = infer_tags(fname, title)
        summary = summarize(body_md)
        notes_data[slug] = {'slug': slug, 'title': title, 'category': category,
                            'tags': tags, 'summary': summary, 'assets': file_assets}
        frontmatter = build_frontmatter(title, fname, category, tags,
                                        project_name, summary, images=file_assets)
        tags_badge = ' '.join(f'`#{t}`' for t in tags)
        related = build_related_links(slug, notes_data)
        full_content = f"""{frontmatter}

# {title}

---

{body_md}
{related}
"""
        wiki_path = os.path.join(wiki_dir, slug + '.md')
        with open(wiki_path, 'w', encoding='utf-8') as f:
            f.write(full_content)
        # slug shifted (duplicate removed)? clean the stale note+assets now
        prev_entry = cache.get(fname)
        if prev_entry:
            prev_wiki_file = prev_entry.get('wiki_file')
            if prev_wiki_file and prev_wiki_file != slug + '.md':
                prev_path = os.path.join(wiki_dir, prev_wiki_file)
                if os.path.isfile(prev_path):
                    os.remove(prev_path)
                    print(f"    ↳ removed stale (renamed): {prev_wiki_file}")
                for asset_rel in prev_entry.get('assets') or []:
                    asset_path = os.path.join(assets_dir, os.path.basename(asset_rel))
                    if os.path.isfile(asset_path):
                        os.remove(asset_path)
        cache[fname] = {'hash': fhash, 'size': fsize,
                        'converted': datetime.now().isoformat(),
                        'wiki_file': slug + '.md',
                        'assets': file_assets}
        cache_changed = True
        converted_count += 1
        print(f"    ✓ Written: {wiki_path} ({len(full_content):,} chars)")

    # Runs even when nothing was converted: an all-unchanged re-run after
    # agent enrichment still refreshes Related links from enriched tags.
    print("\nUpdating cross-links...")
    for slug, data in notes_data.items():
        wiki_path = os.path.join(wiki_dir, slug + '.md')
        if not os.path.isfile(wiki_path):
            continue
        with open(wiki_path, 'r', encoding='utf-8') as f:
            content = f.read()
        related = build_related_links(slug, notes_data)
        if RELATED_MARKER in content:
            content = re.sub(re.escape(RELATED_MARKER) + r'\n## Related\n.*$',
                             related.strip('\n'), content, flags=re.DOTALL)
        else:
            content = content.rstrip() + '\n' + related
        with open(wiki_path, 'w', encoding='utf-8') as f:
            f.write(content)
    print("    ✓ Cross-links updated")

    print("\nBuilding index...")
    with open(os.path.join(wiki_dir, '_Index.md'), 'w', encoding='utf-8') as f:
        f.write(build_index(project_name, notes_data))
    print(f"    ✓ Index: {os.path.join(wiki_dir, '_Index.md')}")

    # WIKI is generated output: remove stale notes/assets even when they are
    # absent from an old or damaged cache.
    current_wiki_files = {d['slug'] + '.md' for d in notes_data.values()}
    current_wiki_files.add('_Index.md')
    current_assets = [asset for data in notes_data.values()
                      for asset in data.get('assets', [])]
    cleanup_generated_outputs(wiki_dir, assets_dir, current_wiki_files,
                              current_assets)
    for src_fname in list(cache.keys()):
        if src_fname not in raw_files_set:  # prune deleted sources
            del cache[src_fname]
            cache_changed = True

    if cache_changed or force:
        save_cache(cache_path, cache)
        print(f"\nCache saved: {cache_path}")
    unchanged_count = max(0, len(raw_files) - converted_count - failed_count)
    print(f"\n✅ Done: {converted_count} converted, "
          f"{unchanged_count} unchanged, "
          f"{len(notes_data)} total notes in WIKI/")
    if failed_count:
        print(f"⚠ {failed_count} file(s) could not be converted.")
    return 1 if failed_count else 0


def vault_config_path():
    return os.path.join(os.path.dirname(os.path.dirname(os.path.abspath(__file__))),
                        'vault-config.json')


def load_vault_config(config_path):
    try:
        with open(config_path, 'r', encoding='utf-8') as f:
            data = json.load(f)
            if isinstance(data, dict) and data.get('vault'):
                return data
    except Exception:
        pass
    return {}


def save_vault_config(config_path, vault):
    os.makedirs(os.path.dirname(config_path), exist_ok=True)
    with open(config_path, 'w', encoding='utf-8') as f:
        json.dump({'vault': vault, 'saved': datetime.now().isoformat()},
                  f, indent=2, ensure_ascii=False)


def _has_raw_subfolder(path):
    try:
        return any(os.path.isdir(os.path.join(path, d, 'RAW'))
                   for d in os.listdir(path))
    except OSError:
        return False


def looks_like_vault(path):
    if not os.path.isdir(path):
        return False
    if os.path.isdir(os.path.join(path, '.obsidian')):
        return True
    return _has_raw_subfolder(path)


def obsidian_registry_vaults():
    candidates = []
    locations = []
    appdata = os.environ.get('APPDATA')
    if appdata:
        locations.append(os.path.join(appdata, 'obsidian', 'obsidian.json'))
    locations += [
        os.path.join(os.path.expanduser('~'), 'Library', 'Application Support',
                     'obsidian', 'obsidian.json'),
        os.path.join(os.path.expanduser('~'), '.config', 'obsidian', 'obsidian.json'),
    ]
    for loc in locations:
        if not os.path.isfile(loc):
            continue
        try:
            with open(loc, 'r', encoding='utf-8') as f:
                data = json.load(f)
            for vinfo in (data.get('vaults') or {}).values():
                p = vinfo.get('path') if isinstance(vinfo, dict) else None
                if p and os.path.isdir(p):
                    candidates.append((os.path.abspath(p), f'Obsidian registry ({loc})'))
        except Exception:
            pass
    return candidates


def drive_scan_vaults():
    candidates = []
    names = ['vault', 'obsidian', 'Obsidian Vault',
             os.path.join('Documents', 'Obsidian Vault'),
             os.path.join('obsidian', 'vault')]
    roots = []
    if os.name == 'nt':
        import string
        for letter in string.ascii_uppercase:
            if letter in ('A', 'B'):
                continue
            d = f'{letter}:\\'
            if os.path.isdir(d):
                roots.append(d)
    home = os.path.expanduser('~')
    roots += [home, os.path.join(home, 'Documents')]
    for root in roots:
        for name in names:
            p = os.path.join(root, name)
            if os.path.isdir(p) and looks_like_vault(p):
                candidates.append((os.path.abspath(p), f'drive scan ({root})'))
    return candidates


def auto_detect_vaults():
    seen = {}
    for p, src in obsidian_registry_vaults() + drive_scan_vaults():
        key = os.path.normcase(p)
        if key not in seen:
            seen[key] = (p, src)
    return list(seen.values())


def discover_projects(vault):
    try:
        return sorted(d for d in os.listdir(vault)
                      if os.path.isdir(os.path.join(vault, d))
                      and os.path.isdir(os.path.join(vault, d, 'RAW')))
    except OSError:
        return []


def discover_subfolders(vault):
    try:
        return sorted(d for d in os.listdir(vault)
                      if os.path.isdir(os.path.join(vault, d))
                      and not d.startswith('.') and d not in ('RAW', 'WIKI'))
    except OSError:
        return []


def move_loose_files_into_raw(subdir, only_if_needed=False):
    raw_dir = os.path.join(subdir, 'RAW')
    already_has_raw = os.path.isdir(raw_dir)
    try:
        entries = sorted(os.listdir(subdir))
    except OSError:
        return False, []
    loose = [e for e in entries
             if os.path.isfile(os.path.join(subdir, e))
             and os.path.splitext(e)[1].lower() in SUPPORTED_EXT]
    if only_if_needed and not already_has_raw and not loose:
        return False, []
    created = False
    if not already_has_raw:
        os.makedirs(raw_dir, exist_ok=True)
        created = True
    moved = []
    import shutil
    for entry in loose:
        try:
            shutil.move(os.path.join(subdir, entry), os.path.join(raw_dir, entry))
            moved.append(entry)
        except OSError as e:
            print(f'  ⚠ Could not move {entry}: {e}')
    return created, moved


def run_vault_subfolders(vault, force):
    subs = discover_subfolders(vault)
    if not subs:
        print(f'No subfolders to process in {vault}')
        return 0
    rc = 0
    for name in subs:
        pdir = os.path.join(vault, name)
        created, moved = move_loose_files_into_raw(pdir, only_if_needed=True)
        if not created and not moved and not os.path.isdir(os.path.join(pdir, 'RAW')):
            continue
        print(f'\n===== {name} =====')
        if created:
            print(f'  📁 RAW/ created in {pdir}')
        if moved:
            print(f'  📦 Moved {len(moved)} file(s) into RAW/: {", ".join(moved)}')
        rc = max(rc, run_project(pdir, force))
    return rc


def vault_child_containing(path, vault):
    try:
        rel = os.path.relpath(path, vault)
    except ValueError:
        return None
    if rel == '.' or rel.startswith('..'):
        return None
    first = rel.split(os.sep)[0]
    if first.startswith('.') or first in ('RAW', 'WIKI'):
        return None
    return first


def parse_wiki_index_from_knowledge(knowledge_path):
    try:
        with open(knowledge_path, 'r', encoding='utf-8') as f:
            content = f.read()
    except OSError:
        return None
    m = re.search(r'^\s*(.+?[\\/]WIKI[\\/]_Index\.md)\s*$', content, re.M | re.I)
    if m:
        return m.group(1).strip()
    m = re.search(r'(\S.*?[\\/]WIKI[\\/]_Index\.md)', content, re.I)
    return m.group(1).strip() if m else None


def last_wiki_update(vault_project_dir):
    cache = load_cache(os.path.join(vault_project_dir, '.owiki-cache.json'))
    stamps = [e.get('converted') for e in cache.values()
              if isinstance(e, dict) and e.get('converted')]
    return max(stamps) if stamps else None


def sync_agent_files_from_knowledge(cwd):
    kp = os.path.join(cwd, 'KNOWLEDGE.md')
    ap = os.path.join(cwd, 'AGENTS.md')
    if not (os.path.isfile(kp) and os.path.isfile(ap)):
        print('Not connected to any WIKI: no KNOWLEDGE.md/AGENTS.md in the current folder.')
        print('Run /owiki <vault-subfolder> from this folder to connect it.')
        return 0
    wiki_index = parse_wiki_index_from_knowledge(kp)
    if not wiki_index:
        print('KNOWLEDGE.md found, but no WIKI/_Index.md path could be parsed —')
        print('cannot determine which vault folder this project is connected to.')
        return 1
    vault_project_dir = os.path.dirname(os.path.dirname(os.path.abspath(wiki_index)))
    project_name = os.path.basename(vault_project_dir)
    last_updated = last_wiki_update(vault_project_dir)
    ensure_knowledge_md(cwd, vault_project_dir, project_name, last_updated=last_updated)
    ap2, action = ensure_agents_md(cwd)
    print(f'  📄 KNOWLEDGE.md → {kp}')
    print(f'  📄 AGENTS.md {action} → {ap2}')
    print(f'  ↳ linked to {wiki_index}'
          + (f'  (last wiki update: {last_updated})' if last_updated else ''))
    return 0


def _ask(prompt, default=None):
    try:
        return input(prompt).strip()
    except EOFError:
        return default if default is not None else ''


def manual_vault_entry():
    while True:
        drive = ''
        if os.name == 'nt':  # drive letters only make sense on Windows
            drive = _ask('  Which drive? (examples: C, D, E) [Enter for home]: ')
            drive = drive.rstrip(':').upper()
        folder = _ask('  Which folder? (examples: vault, obsidian, '
                      'Documents/Obsidian Vault): ')
        if not folder:
            print('  ⚠ A folder name is required — try again.\n')
            continue
        if os.path.isabs(folder):
            path = folder
        elif drive and len(drive) == 1 and drive.isalpha():
            path = drive + ':\\' + folder.strip('\\/')
        else:
            path = os.path.join(os.path.expanduser('~'), folder.strip('/\\'))
        print(f'  → resolved: {path}')
        if not os.path.isdir(path):
            print('  ⚠ This path does not exist.')
            if _ask('  Continue anyway? [y/N] ', 'n').lower() not in ('y', 'yes'):
                continue
        elif not looks_like_vault(path):
            print('  ⚠ This does not look like an Obsidian vault '
                  '(no .obsidian/ folder and no RAW/ subfolder).')
            if _ask('  Continue anyway? [y/N] ', 'n').lower() not in ('y', 'yes'):
                continue
        return path


def run_wizard(config_path):
    """Interactive only — refuses to run without a real TTY so a detected
    vault can never be silently auto-accepted in agent shells."""
    if not sys.stdin.isatty():
        print('ERROR: the setup wizard needs an interactive terminal (stdin is not a TTY).')
        hard_gate_error(config_path)
    print('\nNo vault configured. Auto-detecting your Obsidian vault…')
    candidates = auto_detect_vaults()
    vault = None
    if len(candidates) == 1:
        p, src = candidates[0]
        print(f'✓ Found: {p}   (from {src})')
        if _ask('Is this correct? [Y/n] ', 'y').lower() in ('', 'y', 'yes'):
            vault = p
        else:
            print('OK — enter it manually:\n')
    elif len(candidates) > 1:
        print(f'Multiple Obsidian vaults found ({len(candidates)}):')
        for i, (p, src) in enumerate(candidates, 1):
            print(f'  [{i}] {p}   ({src})')
        ans = _ask(f'Choose one [1-{len(candidates)}], or "n" to enter manually: ', '1')
        if ans.lower() in ('n', 'no'):
            print('OK — enter it manually:\n')
        elif ans.isdigit() and 1 <= int(ans) <= len(candidates):
            vault = candidates[int(ans) - 1][0]
        else:
            vault = candidates[0][0]
    else:
        print('⚠ No Obsidian vault found automatically.')
        print('Enter it manually:\n')
    if vault is None:
        vault = manual_vault_entry()
    save_vault_config(config_path, vault)
    print(f'✓ Saved to {config_path}')
    return vault


def hard_gate_error(config_path, candidates=None):
    script = os.path.abspath(__file__)
    if candidates:
        print('ERROR: no Obsidian vault configured, and several were auto-detected —')
        print('there is no terminal here to ask which one you mean:')
        for p, src in candidates:
            print(f'    - {p}   ({src})')
        print('Fix one of the following, then retry:')
        print('  1. Pass the vault explicitly (this also SAVES the config):')
        print('       python owiki_convert.py --vault "<one of the above>" <project>')
        print('  2. Pin it permanently:         set OWIKI_VAULT_PATH=<vault path>')
        print(f'  3. Choose interactively:       python "{script}" --init   (real terminal)')
        sys.exit(1)
    print('ERROR: no Obsidian vault configured and none could be auto-detected.')
    print('Fix one of the following, then retry:')
    print(f'  1. Run the setup wizard once from a REAL terminal:  python "{script}" --init')
    print(f'     (saves the vault path to {config_path})')
    print('  2. Set the environment variable:  OWIKI_VAULT_PATH=<path to vault>')
    print('  3. Pass the vault explicitly (this also SAVES the config):')
    print('       python owiki_convert.py --vault <vault> <project>')
    print('  4. Use a full project path:       python owiki_convert.py <full project path>')
    sys.exit(1)


def resolve_project_root(cwd, explicit=None):
    return explicit if explicit else cwd


def ensure_knowledge_md(project_root, vault_project_dir, project_name, last_updated=None):
    wiki_index = os.path.join(vault_project_dir, 'WIKI', '_Index.md')
    content = (
        f'# {project_name} — Knowledge Base\n'
        f'\n'
        f'This repository\'s domain knowledge lives in the **Obsidian wiki**:\n'
        f'\n'
        f'    {wiki_index}\n'
        f'\n'
        f'- Read `_Index.md` **first** — it lists every note by category.\n'
        f'- Notes are generated by the `owiki` skill from `RAW/` (read-only\n'
        f'  source input). WIKI notes are extracted by the converter and\n'
        f'  enriched by the active agent model (English translation, summaries,\n'
        f'  tags, figure descriptions, body cleanup). The cache preserves\n'
        f'  enrichment for unchanged files.\n'
        f'- Edit the original files in RAW/ and re-run owiki to update the wiki.\n'
        f'- YAML frontmatter (`title`, `category`, `tags`, `source`) is\n'
        f'  programmatically readable by agents.\n'
    )
    if last_updated:
        content += f'\nLast wiki update: {last_updated}\n'
    path = os.path.join(project_root, 'KNOWLEDGE.md')
    with open(path, 'w', encoding='utf-8') as f:
        f.write(content)
    return path


AGENTS_DOMAIN_RULE = (
    '## Domain Knowledge\n'
    '\n'
    'This project\'s domain knowledge (schemas, business logic, specs) lives\n'
    'in **KNOWLEDGE.md** in this repository. Read it before touching domain\n'
    'logic — it points to the Obsidian wiki that holds the full specs.\n'
)


def ensure_agents_md(project_root):
    path = os.path.join(project_root, 'AGENTS.md')
    if not os.path.isfile(path):
        with open(path, 'w', encoding='utf-8') as f:
            f.write(AGENTS_DOMAIN_RULE)
        return path, 'created'
    with open(path, 'r', encoding='utf-8') as f:
        existing = f.read()
    if '## Domain Knowledge' in existing and 'KNOWLEDGE.md' in existing:
        return path, 'unchanged'
    sep = '\n' if existing.endswith('\n') else '\n\n'
    with open(path, 'a', encoding='utf-8') as f:
        f.write(sep + AGENTS_DOMAIN_RULE)
    return path, 'updated'


def generate_agent_files(project_root, vault_project_dir, project_name, last_updated=None):
    kp = ensure_knowledge_md(project_root, vault_project_dir, project_name,
                             last_updated=last_updated)
    ap, action = ensure_agents_md(project_root)
    print(f'  📄 KNOWLEDGE.md → {kp}')
    print(f'  📄 AGENTS.md {action} → {ap}')
    return True


def main():
    parser = argparse.ArgumentParser(
        description='Convert RAW→WIKI for Obsidian vaults (location-agnostic)',
        epilog='Examples:\n'
               '  owiki_convert.py "<vault>/myproject"           full project path (convert only)\n'
               '  owiki_convert.py myproject                     named subfolder (writes KNOWLEDGE.md/AGENTS.md)\n'
               '  owiki_convert.py --vault "<vault>" myproject   explicit vault + project (saves config)\n'
               '  owiki_convert.py --all | --list | --init | --update',
        formatter_class=argparse.RawDescriptionHelpFormatter)
    parser.add_argument('target', nargs='?', help='project name or full project path')
    parser.add_argument('--vault', help='explicit vault root path (saved to config)')
    parser.add_argument('--list', action='store_true', help='list RAW/-bearing projects')
    parser.add_argument('--all', action='store_true', help='convert every vault subfolder')
    parser.add_argument('--init', action='store_true', help='re-run the setup wizard')
    parser.add_argument('--update', action='store_true',
                        help='re-sync KNOWLEDGE.md/AGENTS.md in cwd from the wiki link')
    parser.add_argument('--force', action='store_true', help='Reconvert all files')
    parser.add_argument('--project-root',
                        help='project root dir for KNOWLEDGE.md/AGENTS.md '
                             '(default: the directory you ran /owiki from)')
    args = parser.parse_args()
    config_path = vault_config_path()
    is_tty = sys.stdin.isatty()
    cwd = os.getcwd()

    looks_full_path = (args.target and
                       (os.path.isabs(args.target) or '/' in args.target or '\\' in args.target))
    if looks_full_path and os.path.isdir(args.target):
        sys.exit(run_project(os.path.abspath(args.target), args.force))

    if args.update:
        sys.exit(sync_agent_files_from_knowledge(cwd))

    vault = None
    if args.vault:
        vault = os.path.abspath(args.vault)
    elif os.environ.get('OWIKI_VAULT_PATH'):
        vault = os.path.abspath(os.environ['OWIKI_VAULT_PATH'])
    else:
        cfg = load_vault_config(config_path)
        if cfg.get('vault'):
            vault = cfg['vault']

    wizard_ran = False
    if args.init:
        if not is_tty:
            print('ERROR: --init needs a terminal (no TTY available).')
            hard_gate_error(config_path)
        vault = run_wizard(config_path)
        wizard_ran = True

    if vault is None:
        if is_tty:
            vault = run_wizard(config_path)
            wizard_ran = True
        else:
            # Agent mode (no TTY): the wizard can't prompt here, but detection
            # itself needs no terminal. Adopt and SAVE a single unambiguous
            # candidate so the config file exists for every future run. When
            # several vaults exist, the one already containing RAW/ project
            # folders wins (that's the owiki vault, not just any Obsidian
            # vault). Zero or still-ambiguous candidates hit the hard gate.
            candidates = auto_detect_vaults()
            pick = None
            if len(candidates) == 1:
                pick = candidates[0]
            else:
                with_raw = [c for c in candidates if _has_raw_subfolder(c[0])]
                if len(with_raw) == 1:
                    pick = with_raw[0]
            if pick:
                vault, src = pick
                print(f'No vault configured — auto-detected: {vault}')
                print(f'  (from {src} — saved to {config_path})')
                save_vault_config(config_path, vault)
            else:
                hard_gate_error(config_path, candidates)

    if not os.path.isdir(vault):
        print(f'ERROR: vault path does not exist: {vault}')
        print('Run with --vault <path>, set OWIKI_VAULT_PATH, or re-run --init.')
        sys.exit(1)
    if args.vault:
        # An explicitly passed vault is a deliberate choice — pin it so every
        # future run (any agent, no flags) finds it via vault-config.json.
        save_vault_config(config_path, vault)
    projects = discover_projects(vault)

    if args.target:
        pdir = os.path.join(vault, args.target)
        if not os.path.isdir(os.path.join(pdir, 'RAW')):
            print(f'  ⚠ No RAW/ at {pdir} — creating it (drop files in and re-run).')
            os.makedirs(os.path.join(pdir, 'RAW'), exist_ok=True)
        rc = run_project(pdir, args.force)
        proj_root = resolve_project_root(cwd, args.project_root)
        generate_agent_files(proj_root, pdir, args.target,
                             last_updated=last_wiki_update(pdir))
        sys.exit(rc)

    if args.list:
        if not projects:
            print(f'No projects (RAW/-bearing subfolders) found in {vault}')
            sys.exit(1)
        print(f'Vault: {vault}')
        print('Projects:')
        for name in projects:
            raw = os.path.join(vault, name, 'RAW')
            n = len([f for f in os.listdir(raw) if os.path.isfile(os.path.join(raw, f))])
            print(f'  - {name}  ({n} raw files)')
        sys.exit(0)

    if args.all:
        print(f'Vault: {vault} — converting all subfolders')
        sys.exit(run_vault_subfolders(vault, args.force))

    if wizard_ran and is_tty:
        if not projects:
            print(f'No projects (RAW/-bearing subfolders) found in {vault}')
        else:
            print(f'Vault: {vault}')
            print('Projects:')
            for i, name in enumerate(projects, 1):
                print(f'  [{i}] {name}')
            if _ask('\nConvert ALL subfolders now? (same as "/owiki --all") [Y/n] ',
                    'y').lower() in ('', 'y', 'yes'):
                run_vault_subfolders(vault, args.force)
            else:
                print('OK — nothing converted. Config saved.')
        sys.exit(0)

    norm_cwd = os.path.normcase(os.path.abspath(cwd))
    norm_vault = os.path.normcase(os.path.abspath(vault))
    if norm_cwd == norm_vault:
        print(f'Vault root: {vault} — converting all subfolders')
        sys.exit(run_vault_subfolders(vault, args.force))
    sub = vault_child_containing(cwd, vault)
    if sub:
        pdir = os.path.join(vault, sub)
        print(f'Inside vault subfolder: {sub} — converting it')
        created, moved = move_loose_files_into_raw(pdir)
        if created:
            print(f'  📁 RAW/ created in {pdir}')
        if moved:
            print(f'  📦 Moved {len(moved)} file(s) into RAW/: {", ".join(moved)}')
        sys.exit(run_project(pdir, args.force))
    print(f'Not inside the vault ({vault}) — nothing to convert.')
    print('  · cd into the vault root       → /owiki converts every subfolder')
    print('  · cd into a vault subfolder    → /owiki converts just it')
    print('  · /owiki <name>                → convert a named subfolder (writes KNOWLEDGE.md/AGENTS.md)')
    print('  · /owiki --all                 → convert all subfolders from anywhere')
    print('  · /owiki --update              → re-sync KNOWLEDGE.md/AGENTS.md in cwd')
    sys.exit(0)


if __name__ == '__main__':
    main()